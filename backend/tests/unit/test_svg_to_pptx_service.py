"""Unit tests for services.svg_to_pptx_service (SVG -> native editable PPTX)."""
import io

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from services.svg_to_pptx_service import (
    SvgToPptxService,
    _parse_color,
    _strip_to_overlay,
    _Scale,
    SLIDE_W_EMU,
    SLIDE_H_EMU,
)

# A page exercising every translator branch: gradient bg, rounded card, plain rect,
# circle, line, multi-run text (per-tspan colour), and a decorative <path> (overlay).
_PAGE = (
    '<svg viewBox="0 0 1280 720" xmlns="http://www.w3.org/2000/svg">'
    '<defs>'
    '<linearGradient id="bg" x1="0" y1="0" x2="1280" y2="720" gradientUnits="userSpaceOnUse">'
    '<stop offset="0" stop-color="#07111F"/><stop offset="1" stop-color="#101827"/>'
    '</linearGradient>'
    '</defs>'
    '<g data-role="shape">'
    '<rect x="0" y="0" width="1280" height="720" fill="url(#bg)"/>'
    '<rect x="40" y="54" width="1200" height="300" rx="34" fill="#0C1628" stroke="#26364F" stroke-width="2"/>'
    '<circle cx="200" cy="200" r="60" fill="#38BDF8"/>'
    '<line x1="70" y1="400" x2="510" y2="400" stroke="#F59E0B" stroke-width="2"/>'
    '</g>'
    '<g data-role="title">'
    '<text x="100" y="160" font-family="Noto Sans CJK SC" font-size="64" '
    'font-weight="bold" fill="#FFFFFF">标题<tspan fill="#F59E0B">高亮</tspan></text>'
    '<text x="640" y="260" text-anchor="middle" font-family="Noto Sans CJK SC" '
    'font-size="32" fill="#9CA3AF">居中副标题</text>'
    '</g>'
    '<g data-role="decoration">'
    '<path d="M0 610 C180 540 310 590 480 515 L1280 720 L0 720 Z" fill="#0B1220"/>'
    '</g>'
    '</svg>'
)


def _open(data: bytes) -> Presentation:
    return Presentation(io.BytesIO(data))


@pytest.mark.unit
def test_slide_size_is_exact_16_9():
    data, _ = SvgToPptxService.build_from_strings([_PAGE])
    prs = _open(data)
    assert prs.slide_width == SLIDE_W_EMU
    assert prs.slide_height == SLIDE_H_EMU


@pytest.mark.unit
def test_raster_mode_is_default_bg_plus_editable_text():
    data, stats = SvgToPptxService.build_from_strings([_PAGE])
    assert stats["mode"] == "raster"
    prs = _open(data)
    slide = prs.slides[0]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    texts = [s.text_frame.text for s in slide.shapes
             if s.has_text_frame and s.text_frame.text.strip()]

    # one full-page background raster + the editable text on top, no native shapes.
    assert len(pictures) == 1
    assert stats["backgrounds"] == 1
    assert stats["shapes"] == 0
    assert "标题高亮" in texts
    assert "居中副标题" in texts
    # background covers the whole slide
    bg = pictures[0]
    assert bg.left == 0 and bg.top == 0
    assert bg.width == SLIDE_W_EMU and bg.height == SLIDE_H_EMU


@pytest.mark.unit
def test_native_shapes_and_text_are_editable():
    data, stats = SvgToPptxService.build_from_strings([_PAGE], mode="native")
    prs = _open(data)
    slide = prs.slides[0]

    autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    texts = [s.text_frame.text for s in slide.shapes
             if s.has_text_frame and s.text_frame.text.strip()]

    # 2 rects + 1 circle = 3 autoshapes; 2 text boxes; exactly 1 decorative overlay.
    assert len(autoshapes) == 3
    assert "标题高亮" in texts
    assert "居中副标题" in texts
    assert len(pictures) == 1
    assert stats["slides"] == 1 and stats["overlays"] == 1


@pytest.mark.unit
@pytest.mark.parametrize("mode", ["raster", "native"])
def test_per_tspan_colour_becomes_separate_runs(mode):
    data, _ = SvgToPptxService.build_from_strings([_PAGE], mode=mode)
    prs = _open(data)
    runs = None
    for s in prs.slides[0].shapes:
        if s.has_text_frame and s.text_frame.text == "标题高亮":
            runs = s.text_frame.paragraphs[0].runs
            break
    assert runs is not None
    assert [r.text for r in runs] == ["标题", "高亮"]
    assert str(runs[0].font.color.rgb) == "FFFFFF"
    assert str(runs[1].font.color.rgb) == "F59E0B"
    assert runs[0].font.bold is True


@pytest.mark.unit
@pytest.mark.parametrize("mode", ["raster", "native"])
def test_multiple_pages_become_multiple_slides(mode):
    data, stats = SvgToPptxService.build_from_strings([_PAGE, _PAGE, _PAGE], mode=mode)
    prs = _open(data)
    assert len(prs.slides) == 3
    assert stats["slides"] == 3


@pytest.mark.unit
def test_unknown_mode_raises():
    with pytest.raises(ValueError):
        SvgToPptxService.build_from_strings([_PAGE], mode="bogus")


@pytest.mark.unit
def test_empty_input_raises():
    with pytest.raises(ValueError):
        SvgToPptxService.build_from_strings([])


@pytest.mark.unit
@pytest.mark.parametrize("paint,expected", [
    ("#fff", "FFFFFF"),
    ("#F59E0B", "F59E0B"),
    ("none", None),
    ("url(#bg)", "GRAD"),
    ("rgb(255, 0, 0)", "FF0000"),
    ("white", "FFFFFF"),
    (None, "MISSING"),
    ("notacolor", "MISSING"),
])
def test_parse_color(paint, expected):
    assert _parse_color(paint) == expected


@pytest.mark.unit
def test_strip_to_overlay_keeps_paths_drops_primitives():
    stripped = _strip_to_overlay(_PAGE)
    assert "<path" in stripped
    # native primitives must be gone from the overlay layer
    assert "<rect" not in stripped
    assert "<circle" not in stripped
    assert "<text" not in stripped
    # gradient defs are preserved (paths may reference them)
    assert "linearGradient" in stripped


@pytest.mark.unit
def test_scale_maps_contract_viewbox_to_emu():
    s = _Scale(1280.0, 720.0)
    # 1280px canvas -> 9525 EMU/px (96 dpi); font px -> 0.75 pt.
    assert int(s.emu_x) == 9525
    assert round(s.pt, 4) == 0.75
