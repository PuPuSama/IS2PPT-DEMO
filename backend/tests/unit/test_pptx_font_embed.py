"""Unit tests for services.pptx_font_embed (embed a subsetted font into a .pptx)."""
import io
import os
import zipfile

import pytest
from pptx import Presentation

from services.pptx_font_embed import embed_font_into_pptx_bytes

_FONT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "fonts", "NotoSansSC-Regular.ttf",
)
_FAMILY = "Noto Sans CJK SC"


def _blank_pptx_bytes() -> bytes:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.mark.unit
def test_embeds_font_part_and_declarations():
    data = embed_font_into_pptx_bytes(_blank_pptx_bytes(), "标题避雷ABC", _FAMILY, _FONT)
    z = zipfile.ZipFile(io.BytesIO(data))
    names = z.namelist()

    assert "ppt/fonts/font1.fntdata" in names
    pres = z.read("ppt/presentation.xml").decode("utf-8")
    assert "embeddedFontLst" in pres
    assert 'embedTrueTypeFonts="1"' in pres
    assert _FAMILY in pres
    # all four weight/style slots must be declared (decks use bold titles), else
    # PowerPoint can't match a bold run and substitutes a system font
    for slot in ("regular", "bold", "italic", "boldItalic"):
        assert f"<p:{slot} " in pres
    ct = z.read("[Content_Types].xml").decode("utf-8")
    assert "fntdata" in ct
    rels = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    assert "fonts/font1.fntdata" in rels


@pytest.mark.unit
def test_result_is_still_a_valid_pptx():
    data = embed_font_into_pptx_bytes(_blank_pptx_bytes(), "测试", _FAMILY, _FONT)
    # python-pptx must still be able to open the rewritten package
    Presentation(io.BytesIO(data))


@pytest.mark.unit
def test_subset_is_small_but_covers_used_chars():
    data = embed_font_into_pptx_bytes(_blank_pptx_bytes(), "避雷指南", _FAMILY, _FONT)
    z = zipfile.ZipFile(io.BytesIO(data))
    font_bytes = z.read("ppt/fonts/font1.fntdata")
    # subset must be a tiny fraction of the full ~10MB+ CJK font
    assert len(font_bytes) < 1_000_000

    from fontTools.ttLib import TTFont
    f = TTFont(io.BytesIO(font_bytes))
    assert f["name"].getDebugName(1) == _FAMILY  # family preserved -> PowerPoint matches
    cmap = f.getBestCmap()
    for ch in "避雷指南A1":
        assert ord(ch) in cmap


@pytest.mark.unit
def test_best_effort_returns_input_on_bad_font():
    original = _blank_pptx_bytes()
    out = embed_font_into_pptx_bytes(original, "x", _FAMILY, "/no/such/font.ttf")
    assert out == original  # never raise / corrupt the export


@pytest.mark.unit
def test_no_chars_is_noop():
    original = _blank_pptx_bytes()
    assert embed_font_into_pptx_bytes(original, "", _FAMILY, _FONT) == original
