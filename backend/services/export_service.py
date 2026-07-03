"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import math
import os
import json
import logging
import random
import re
import tempfile
import base64
import hashlib
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from textwrap import dedent
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from PIL import Image
import io
import tempfile
import img2pdf
import fitz  # PyMuPDF
from utils.pptx_math import latex_to_display_text, looks_like_latex_math
logger = logging.getLogger(__name__)


class ExportError(Exception):
    """
    导出过程中的错误异常

    当 fail_fast=True 时，任何导出错误都会抛出此异常，
    包含详细的错误信息和帮助提示。
    """
    def __init__(self, message: str, error_type: str = 'unknown', details: Dict[str, Any] = None, help_text: str = None):
        """
        Args:
            message: 错误消息
            error_type: 错误类型 (style_extraction, text_render, image_add, inpaint, config, service)
            details: 详细错误信息
            help_text: 帮助提示文本
        """
        super().__init__(message)
        self.message = message
        self.error_type = error_type
        self.details = details or {}
        self.help_text = help_text or self._get_default_help_text(error_type)

    def _get_default_help_text(self, error_type: str) -> str:
        """根据错误类型返回默认帮助提示"""
        help_texts = {
            'style_extraction': '样式提取失败可能是由于百度OCR API配置问题。请检查「项目设置 -> 导出设置」中的配置，或尝试切换到「MinerU提取」方法。',
            'text_render': '文本渲染失败可能是由于字体或编码问题。请检查页面内容是否包含特殊字符。',
            'image_add': '图片添加失败可能是由于图片文件损坏或路径错误。请尝试重新生成该页面的图片。',
            'inpaint': '背景修复失败可能是由于API配置问题。请检查「项目设置 -> 导出设置」中的背景图获取方法配置。',
            'config': '配置错误。请检查「项目设置 -> 导出设置」中的相关配置。',
            'service': '服务不可用。请稍后重试或联系管理员。',
        }
        return help_texts.get(error_type, '如果问题持续出现，可以在「项目设置 -> 导出设置」中开启「返回半成品」选项以跳过错误继续导出。')

    def to_dict(self) -> Dict[str, Any]:
        """转换为字典格式"""
        return {
            'message': self.message,
            'error_type': self.error_type,
            'details': self.details,
            'help_text': self.help_text
        }


@dataclass
class ExportWarnings:
    """
    导出过程中收集的警告信息
    
    用于追踪哪些操作没有按预期执行，并反馈给前端
    """
    # 样式提取失败的元素
    style_extraction_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 文本渲染失败的元素
    text_render_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 图片添加失败
    image_add_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # JSON 解析失败（重试后仍失败）
    json_parse_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 其他警告
    other_warnings: List[str] = field(default_factory=list)
    
    def add_style_extraction_failed(self, element_id: str, reason: str):
        """记录样式提取失败"""
        self.style_extraction_failed.append({
            'element_id': element_id,
            'reason': reason
        })
    
    def add_text_render_failed(self, text: str, reason: str):
        """记录文本渲染失败"""
        self.text_render_failed.append({
            'text': text[:50] + '...' if len(text) > 50 else text,
            'reason': reason
        })
    
    def add_image_failed(self, path: str, reason: str):
        """记录图片添加失败"""
        self.image_add_failed.append({
            'path': path,
            'reason': reason
        })
    
    def add_json_parse_failed(self, context: str, reason: str):
        """记录 JSON 解析失败"""
        self.json_parse_failed.append({
            'context': context,
            'reason': reason
        })
    
    def add_warning(self, message: str):
        """添加其他警告"""
        self.other_warnings.append(message)
    
    def has_warnings(self) -> bool:
        """是否有警告"""
        return bool(
            self.style_extraction_failed or 
            self.text_render_failed or 
            self.image_add_failed or
            self.json_parse_failed or
            self.other_warnings
        )
    
    def to_summary(self) -> List[str]:
        """生成警告摘要（适合前端展示）"""
        summary = []
        
        if self.style_extraction_failed:
            summary.append(f"⚠️ {len(self.style_extraction_failed)} 个文本元素样式提取失败")
        
        if self.text_render_failed:
            summary.append(f"⚠️ {len(self.text_render_failed)} 个文本元素渲染失败")
        
        if self.image_add_failed:
            summary.append(f"⚠️ {len(self.image_add_failed)} 张图片添加失败")
        
        if self.json_parse_failed:
            summary.append(f"⚠️ {len(self.json_parse_failed)} 次 AI 响应解析失败")
        
        for warning in self.other_warnings[:5]:  # 最多显示5条其他警告
            summary.append(f"⚠️ {warning}")
        
        if len(self.other_warnings) > 5:
            summary.append(f"  ...还有 {len(self.other_warnings) - 5} 条其他警告")
        
        return summary
    
    def to_dict(self) -> Dict[str, Any]:
        """转换为字典（详细信息）"""
        return {
            'style_extraction_failed': self.style_extraction_failed,
            'text_render_failed': self.text_render_failed,
            'image_add_failed': self.image_add_failed,
            'json_parse_failed': self.json_parse_failed,
            'other_warnings': self.other_warnings,
            'total_warnings': (
                len(self.style_extraction_failed) + 
                len(self.text_render_failed) + 
                len(self.image_add_failed) +
                len(self.json_parse_failed) +
                len(self.other_warnings)
            )
        }


def _get_page_size_inches(aspect_ratio: str = '16:9', base: float = 10.0) -> Tuple[float, float]:
    """Return (width, height) in inches for a given aspect ratio string."""
    try:
        w, h = (float(x) for x in aspect_ratio.split(':'))
        if not (math.isfinite(w) and math.isfinite(h) and w > 0 and h > 0):
            raise ValueError(f"invalid dimensions: {w}:{h}")
    except (ValueError, AttributeError) as e:
        logger.warning(f"Invalid aspect ratio '{aspect_ratio}', falling back to 16:9: {e}")
        w, h = 16.0, 9.0
    if w >= h:
        return base, base * h / w
    else:
        return base * w / h, base


class ExportService:
    """Service for exporting presentations"""

    PPTX_TRANSITION_EFFECTS = {
        'fade',
        'page_turn',
        'push',
        'wipe',
        'split',
        'blinds',
        'checker',
        'wheel',
    }

    @staticmethod
    def _apply_slide_transition(slide, effect: str) -> None:
        """Add a PowerPoint slide transition node to a slide XML element."""
        transition = OxmlElement('p:transition')
        transition.set('spd', 'med')

        if effect == 'fade':
            transition.append(OxmlElement('p:fade'))
        elif effect == 'page_turn':
            cover = OxmlElement('p:cover')
            cover.set('dir', 'l')
            transition.append(cover)
        elif effect == 'push':
            push = OxmlElement('p:push')
            push.set('dir', 'l')
            transition.append(push)
        elif effect == 'wipe':
            wipe = OxmlElement('p:wipe')
            wipe.set('dir', 'l')
            transition.append(wipe)
        elif effect == 'split':
            split = OxmlElement('p:split')
            split.set('orient', 'horz')
            split.set('dir', 'out')
            transition.append(split)
        elif effect == 'blinds':
            blinds = OxmlElement('p:blinds')
            blinds.set('dir', 'vert')
            transition.append(blinds)
        elif effect == 'checker':
            checker = OxmlElement('p:checker')
            checker.set('dir', 'horz')
            transition.append(checker)
        elif effect == 'wheel':
            wheel = OxmlElement('p:wheel')
            wheel.set('spokes', '1')
            transition.append(wheel)
        else:
            return

        slide_element = slide._element
        existing = slide_element.find(qn('p:transition'))
        if existing is not None:
            slide_element.remove(existing)

        clr_map_ovr = slide_element.find(qn('p:clrMapOvr'))
        if clr_map_ovr is not None:
            insert_at = slide_element.index(clr_map_ovr) + 1
        else:
            c_sld = slide_element.find(qn('p:cSld'))
            insert_at = slide_element.index(c_sld) + 1 if c_sld is not None else 0

        slide_element.insert(insert_at, transition)

    @staticmethod
    def create_pptx_from_images(
        image_paths: List[str],
        output_file: str = None,
        aspect_ratio: str = '16:9',
        transition_effects: Optional[List[str]] = None,
    ) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set author/date metadata for exported PPTX
        try:
            core = prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "banana-slides"
            core.last_modified_by = "banana-slides"
            core.created = now
            core.modified = now
            core.last_printed = None
        except Exception as e:
            logger.warning(f"Failed to set core properties: {e}")
        
        # Set slide dimensions based on aspect ratio
        page_w, page_h = _get_page_size_inches(aspect_ratio)
        prs.slide_width = Inches(page_w)
        prs.slide_height = Inches(page_h)
        
        valid_transition_effects = [
            effect for effect in (transition_effects or [])
            if effect in ExportService.PPTX_TRANSITION_EFFECTS
        ]
        transition_effect_queue: List[str] = []

        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            if valid_transition_effects:
                if not transition_effect_queue:
                    transition_effect_queue = valid_transition_effects[:]
                    random.shuffle(transition_effect_queue)
                ExportService._apply_slide_transition(
                    slide,
                    transition_effect_queue.pop(),
                )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None, aspect_ratio: str = '16:9') -> Optional[bytes]:
        """
        Create PDF file from image paths using img2pdf (low memory usage)

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        # Validate images exist and log warnings for missing files
        valid_paths = []
        for p in image_paths:
            if os.path.exists(p):
                valid_paths.append(p)
            else:
                logger.warning(f"Image not found and will be skipped for PDF export: {p}")

        if not valid_paths:
            raise ValueError("No valid images found for PDF export")

        try:
            logger.info(f"Using img2pdf for PDF export ({len(valid_paths)} pages, low memory mode)")

            page_w, page_h = _get_page_size_inches(aspect_ratio)
            layout_fun = img2pdf.get_layout_fun(
                pagesize=(img2pdf.in_to_pt(page_w), img2pdf.in_to_pt(page_h)),
                fit=img2pdf.FitMode.fill,
            )

            # Convert images to PDF
            pdf_bytes = img2pdf.convert(valid_paths, layout_fun=layout_fun)

            # Add metadata
            pdf_bytes = ExportService._add_pdf_metadata(pdf_bytes)

            if output_file:
                with open(output_file, "wb") as f:
                    f.write(pdf_bytes)
                return None
            else:
                return pdf_bytes
        except (img2pdf.ImageOpenError, ValueError, IOError) as e:
            logger.warning(f"img2pdf conversion failed: {e}. Falling back to Pillow (high memory usage).")
            return ExportService.create_pdf_from_images_pillow(valid_paths, output_file, aspect_ratio)

    @staticmethod
    def _add_pdf_metadata(pdf_bytes: bytes) -> bytes:
        """Add author metadata to PDF (including XMP for Windows compatibility)"""
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")

            doc.set_metadata({
                "author": "banana-slides",
                "producer": "banana-slides",
                "creator": "banana-slides"
            })

            now = datetime.now(timezone.utc)
            iso_time = now.isoformat()

            content_hash = hashlib.md5(pdf_bytes[:1024]).hexdigest()

            xmp = dedent(f'''\
                <?xpacket begin="" id="W5M0MpCehiHzreSzNTczkc9d"?>
                <x:xmpmeta xmlns:x="adobe:ns:meta/">
                  <rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#">
                    <rdf:Description rdf:about="" xmlns:dc="http://purl.org/dc/elements/1.1/">
                      <dc:creator><rdf:Seq><rdf:li>banana-slides</rdf:li></rdf:Seq></dc:creator>
                    </rdf:Description>
                    <rdf:Description rdf:about="" xmlns:pdf="http://ns.adobe.com/pdf/1.3/">
                      <pdf:Producer>banana-slides</pdf:Producer>
                    </rdf:Description>
                    <rdf:Description rdf:about="" xmlns:xmp="http://ns.adobe.com/xap/1.0/">
                      <xmp:CreatorTool>banana-slides</xmp:CreatorTool>
                      <xmp:CreateDate>{iso_time}</xmp:CreateDate>
                      <xmp:MetadataDate>{iso_time}</xmp:MetadataDate>
                    </rdf:Description>
                    <rdf:Description rdf:about="" xmlns:xmpMM="http://ns.adobe.com/xap/1.0/mm/">
                      <xmpMM:DocumentID>uuid:{content_hash}</xmpMM:DocumentID>
                    </rdf:Description>
                  </rdf:RDF>
                </x:xmpmeta>
                <?xpacket end="w"?>''')
            doc.set_xml_metadata(xmp)

            return doc.tobytes()
        except Exception as e:
            logger.warning(f"Failed to add PDF metadata: {e}")
            return pdf_bytes

    @staticmethod
    def create_pdf_from_images_pillow(image_paths: List[str], output_file: str = None, aspect_ratio: str = '16:9') -> Optional[bytes]:
        """
        Create PDF file from image paths using Pillow (original method)

        Note: This method loads all images into memory at once.
        For large projects (50+ pages with 20MB/page), use create_pdf_from_images instead.

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        images = []
        page_w, page_h = _get_page_size_inches(aspect_ratio)

        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue

            img = Image.open(image_path)

            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')

            # Set DPI so PDF page matches target dimensions
            img.info['dpi'] = (img.width / page_w, img.height / page_h)

            images.append(img)

        if not images:
            raise ValueError("No valid images found for PDF export")

        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return ExportService._add_pdf_metadata(pdf_bytes.getvalue())
