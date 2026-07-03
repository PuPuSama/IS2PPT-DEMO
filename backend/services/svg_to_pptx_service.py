"""
SVG → editable PPTX translator (svg generation mode, stage 3).

Turns the SVG pages produced by the 'svg' generation mode into NATIVE, click-to-edit
PowerPoint shapes — no MinerU, no OCR, no inpainting. Each SVG primitive maps to a
real pptx object:

  <rect>            -> rectangle / rounded-rectangle (fill, line, corner radius, gradient)
  <circle><ellipse> -> oval
  <line>            -> straight connector
  <text>/<tspan>    -> text box (font, size, weight, colour, alignment; per-tspan
                       colour becomes a separate run; ancestor <g> attrs inherited)

Everything that can't be a clean editable primitive — decorative ``<path>`` (icons,
arrows, dividers) and elements inside faint ``<g opacity=...>`` decoration groups — is
composited into ONE transparent, full-canvas PNG overlaid on top of the slide. That
overlay is rendered with **resvg** (the same renderer the svg generation mode already
uses), so there is no Chrome/headless-browser dependency.

Coordinate system: the SVG contract locks ``viewBox="0 0 1280 720"`` (96 dpi → a
13.333 × 7.5 in 16:9 slide, 9525 EMU/px), but the viewBox is read at runtime and the
EMU/pt scaling derived from it, so any 16:9-ish canvas still maps correctly.

Fonts: text runs keep the SVG's declared ``font-family`` (the bundled CJK face is
"Noto Sans CJK SC"); the resvg overlay loads only ``backend/fonts/`` so the rasterised
decorations match the editable text.
"""
import os
import re
import logging
from io import BytesIO
from typing import List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

from services.svg_render_service import render_svg_to_png

logger = logging.getLogger(__name__)

SVG_NS = "http://www.w3.org/2000/svg"

# 16:9 slide canvas (EMU). 12192000 EMU = 13.333 in, 6858000 EMU = 7.5 in.
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
EMU_PER_PT = 12700

# Default logical canvas if a viewBox is missing/unparsable.
DEFAULT_VBW = 1280.0
DEFAULT_VBH = 720.0

# Baseline offset above the text box top, in em units (calibrated against resvg).
BASELINE_EM = 0.86

# Fallback CJK family if an SVG omits font-family.
# Fallback when an SVG omits font-family. Text runs otherwise pass the SVG's declared
# font-family ("Noto Sans CJK SC") through to PowerPoint; that font is also EMBEDDED in
# the exported pptx (see pptx_font_embed) so it matches the preview on any machine.
DEFAULT_FONT_FAMILY = "Noto Sans CJK SC"

# Elements translated to native shapes (everything else -> overlay).
SHAPELIKE = ("rect", "circle", "ellipse", "line", "polyline", "polygon", "text")

_NAMED_COLORS = {
    "white": "FFFFFF", "black": "000000", "red": "FF0000",
    "green": "008000", "blue": "0000FF", "gray": "808080", "grey": "808080",
}


# --------------------------------------------------------------------------- #
# small helpers
# --------------------------------------------------------------------------- #
def _localname(el) -> str:
    """Tag local name, '' for comment / PI nodes."""
    tag = el.tag
    if not isinstance(tag, str):
        return ""
    return etree.QName(tag).localname


def _parse_color(c: Optional[str]) -> Optional[str]:
    """SVG paint string -> 'RRGGBB' | None (no paint) | 'GRAD' (url ref) | 'MISSING'."""
    if c is None:
        return "MISSING"
    c = c.strip()
    if c == "" or c == "none":
        return None
    if c.startswith("url("):
        return "GRAD"
    if c.startswith("#"):
        h = c[1:]
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        if len(h) == 6:
            return h.upper()
        return "MISSING"
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", c)
    if m:
        return "{:02X}{:02X}{:02X}".format(*(int(v) for v in m.groups()))
    return _NAMED_COLORS.get(c.lower(), "MISSING")


def _inherited(el, attr: str, default: Optional[str] = None) -> Optional[str]:
    """Resolve a presentation attribute, walking up ancestors (closest wins)."""
    cur = el
    while cur is not None and etree.iselement(cur):
        v = cur.get(attr)
        if v is not None:
            return v
        cur = cur.getparent()
    return default


def _in_faint_group(el) -> bool:
    """True if any ancestor <g> carries an opacity attribute (faint decoration)."""
    cur = el.getparent()
    while cur is not None and etree.iselement(cur):
        if _localname(cur) == "g" and cur.get("opacity") is not None:
            return True
        cur = cur.getparent()
    return False


def _gradient_stops(root, gid: str):
    """[(offset, 'RRGGBB'), ...] for a linear/radial gradient id, or None."""
    for kind in ("linearGradient", "radialGradient"):
        for grad in root.iter("{%s}%s" % (SVG_NS, kind)):
            if grad.get("id") == gid:
                stops = []
                for s in grad:
                    if _localname(s) == "stop":
                        col = _parse_color(s.get("stop-color"))
                        if col in (None, "GRAD", "MISSING"):
                            col = "808080"
                        stops.append((float(s.get("offset", 0) or 0), col))
                return stops or None
    return None


class _Scale:
    """Maps SVG user units to EMU / pt for a given viewBox."""

    def __init__(self, vbw: float, vbh: float):
        self.vbw = vbw or DEFAULT_VBW
        self.vbh = vbh or DEFAULT_VBH
        self.emu_x = SLIDE_W_EMU / self.vbw
        self.emu_y = SLIDE_H_EMU / self.vbh
        # font size is a length; use the X scale (uniform for a 16:9 viewBox).
        self.pt = self.emu_x / EMU_PER_PT

    def ex(self, v) -> Emu:
        return Emu(int(round(float(v) * self.emu_x)))

    def ey(self, v) -> Emu:
        return Emu(int(round(float(v) * self.emu_y)))


# --------------------------------------------------------------------------- #
# styling
# --------------------------------------------------------------------------- #
def _apply_solid_or_none(shape, paint: Optional[str]):
    if paint is None:
        shape.fill.background()
    elif paint == "MISSING":
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("000000")
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(paint)


def _apply_gradient(shape, root, fill_ref: str):
    m = re.search(r"url\(#([^)]+)\)", fill_ref or "")
    stops = _gradient_stops(root, m.group(1)) if m else None
    if not stops:
        stops = [(0.0, "0B2A4A"), (1.0, "11406F")]
    try:
        shape.fill.gradient()
        gss = shape.fill.gradient_stops
        for i, (off, col) in enumerate(stops[:len(gss)]):
            gss[i].position = max(0.0, min(1.0, off))
            gss[i].color.rgb = RGBColor.from_string(col)
        shape.fill.gradient_angle = 45.0
    except Exception:
        # python-pptx only exposes 2-stop gradients on some autoshapes; fall back.
        _apply_solid_or_none(shape, stops[0][1])


def _apply_line(shape, stroke: Optional[str], width, scale: _Scale):
    sc = _parse_color(stroke) if stroke is not None else None
    if sc and sc not in ("GRAD", "MISSING"):
        shape.line.color.rgb = RGBColor.from_string(sc)
        shape.line.width = scale.ex(float(width or 1))
    else:
        shape.line.fill.background()


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _set_cjk_typeface(run, name: str):
    """Force the east-asian / complex-script typeface (python-pptx only sets latin)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)


def _set_letter_spacing(run, px: float, scale: _Scale):
    # OOXML 'spc' is in 1/100 pt.
    run._r.get_or_add_rPr().set("spc", str(int(px * scale.pt * 100)))


def _estimate_text_width(runs, font_size: float, spacing: float) -> float:
    """Rough rendered width in user units (CJK ~1em, latin ~0.55em)."""
    n = 0
    total = 0.0
    for txt, _, _ in runs:
        for ch in txt:
            n += 1
            total += font_size * (0.55 if ord(ch) < 0x2E80 else 1.0)
    return total + spacing * max(0, n - 1)


# --------------------------------------------------------------------------- #
# element translators
# --------------------------------------------------------------------------- #
def _add_rect(shapes, el, root, scale: _Scale):
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    w = float(el.get("width", 0)); h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        return
    rx = el.get("rx")
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(shp, scale.ex(x), scale.ey(y), scale.ex(w), scale.ey(h))
    if rx:
        try:
            sp.adjustments[0] = max(0.0, min(0.5, float(rx) / min(w, h)))
        except Exception:
            pass
    fill_ref = _inherited(el, "fill")
    fc = _parse_color(fill_ref)
    if fc == "GRAD":
        _apply_gradient(sp, root, fill_ref)
    else:
        _apply_solid_or_none(sp, fc)
    _apply_line(sp, _inherited(el, "stroke"), _inherited(el, "stroke-width"), scale)
    _no_shadow(sp)


def _add_oval(shapes, el, scale: _Scale):
    name = _localname(el)
    if name == "circle":
        cx = float(el.get("cx", 0)); cy = float(el.get("cy", 0)); r = float(el.get("r", 0))
        rx = ry = r
    else:  # ellipse
        cx = float(el.get("cx", 0)); cy = float(el.get("cy", 0))
        rx = float(el.get("rx", 0)); ry = float(el.get("ry", 0))
    if rx <= 0 or ry <= 0:
        return
    sp = shapes.add_shape(MSO_SHAPE.OVAL,
                          scale.ex(cx - rx), scale.ey(cy - ry),
                          scale.ex(2 * rx), scale.ey(2 * ry))
    fc = _parse_color(_inherited(el, "fill"))
    if fc == "GRAD":
        _apply_solid_or_none(sp, "MISSING")
    else:
        _apply_solid_or_none(sp, fc)
    _apply_line(sp, _inherited(el, "stroke"), _inherited(el, "stroke-width"), scale)
    _no_shadow(sp)


def _add_line(shapes, el, scale: _Scale):
    cn = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        scale.ex(el.get("x1", 0)), scale.ey(el.get("y1", 0)),
        scale.ex(el.get("x2", 0)), scale.ey(el.get("y2", 0)),
    )
    sc = _parse_color(_inherited(el, "stroke"))
    if sc and sc not in ("GRAD", "MISSING"):
        cn.line.color.rgb = RGBColor.from_string(sc)
    cn.line.width = scale.ex(float(_inherited(el, "stroke-width") or 1))
    _no_shadow(cn)


def _add_text(shapes, el, scale: _Scale):
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    fs = float(_inherited(el, "font-size", 16))
    anchor = _inherited(el, "text-anchor", "start")
    spacing_raw = _inherited(el, "letter-spacing")
    spacing = float(spacing_raw) if spacing_raw else 0.0
    family = _inherited(el, "font-family", DEFAULT_FONT_FAMILY) or DEFAULT_FONT_FAMILY
    family = family.split(",")[0].strip().strip("'\"") or DEFAULT_FONT_FAMILY

    base_fill = _parse_color(_inherited(el, "fill"))
    if base_fill in (None, "GRAD", "MISSING"):
        base_fill = "000000"
    base_bold = (_inherited(el, "font-weight") in ("bold", "600", "700", "800", "900"))

    # Collect runs: element text, then each <tspan> (own colour/weight) + its tail.
    runs: List[Tuple[str, str, bool]] = []
    if el.text and el.text.strip():
        runs.append((el.text, base_fill, base_bold))
    for child in el:
        if _localname(child) == "tspan":
            tf = _parse_color(child.get("fill"))
            if tf in (None, "GRAD", "MISSING"):
                tf = base_fill
            tb = (child.get("font-weight") in ("bold", "600", "700", "800", "900")) or base_bold
            if child.text:
                runs.append((child.text, tf, tb))
        if child.tail and child.tail.strip():
            runs.append((child.tail, base_fill, base_bold))
    if not runs:
        return

    width = max(_estimate_text_width(runs, fs, spacing), 4.0)
    left = {"start": x, "middle": x - width / 2, "end": x - width}.get(anchor, x)
    top = y - fs * BASELINE_EM

    box = shapes.add_textbox(scale.ex(left), scale.ey(top),
                             scale.ex(width), scale.ey(fs * 1.45))
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER,
                   "end": PP_ALIGN.RIGHT}.get(anchor, PP_ALIGN.LEFT)
    for txt, col, bold in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(fs * scale.pt)
        r.font.bold = bold
        r.font.name = family
        r.font.color.rgb = RGBColor.from_string(col)
        _set_cjk_typeface(r, family)
        if spacing:
            _set_letter_spacing(r, spacing, scale)


# --------------------------------------------------------------------------- #
# decorative overlay (resvg)
# --------------------------------------------------------------------------- #
def _strip_to_overlay(svg_text: str) -> str:
    """Keep only <path> + faint-group decorations (and <defs>); drop native primitives."""
    root = etree.fromstring(svg_text.encode("utf-8"))
    for el in list(root.iter()):
        if _localname(el) in SHAPELIKE and not _in_faint_group(el):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
    return etree.tostring(root, encoding="unicode")


# --------------------------------------------------------------------------- #
# fidelity (raster background) mode
# --------------------------------------------------------------------------- #
def _editable_texts(root):
    """The <text> elements that become native editable boxes (non-faint, has text)."""
    out = []
    for el in root.iter():
        if _localname(el) == "text" and not _in_faint_group(el):
            has = (el.text and el.text.strip()) or any(
                (c.text and c.text.strip()) or (c.tail and c.tail.strip()) for c in el
            )
            if has:
                out.append(el)
    return out


def _strip_editable_text(svg_text: str) -> str:
    """Remove the <text> elements that will be re-added as native boxes.

    Everything else (shapes, gradients, paths, icons, opacity, faint decorations,
    AND any faint/decorative text) stays so the rasterised background is pixel-
    identical to the web preview except where the editable text used to be.
    """
    root = etree.fromstring(svg_text.encode("utf-8"))
    for el in list(root.iter()):
        if _localname(el) == "text" and not _in_faint_group(el):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
    return etree.tostring(root, encoding="unicode")


def _render_background(svg_text: str, bg_width: int) -> Optional[BytesIO]:
    """Render the full page minus editable text to a PNG (the preview's own pixels)."""
    stripped = _strip_editable_text(svg_text)
    try:
        img = render_svg_to_png(stripped, width=int(bg_width))
    except Exception as e:
        logger.warning(f"SVG background render failed: {e}")
        return None
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _render_overlay(svg_text: str, vbw: float, overlay_scale: int) -> Optional[BytesIO]:
    """Render the decorative-only SVG to a transparent PNG; None if nothing to draw."""
    stripped = _strip_to_overlay(svg_text)
    # Anything left worth rasterising? (a bare <svg><defs/></svg> is not)
    probe = etree.fromstring(stripped.encode("utf-8"))
    if not any(_localname(el) in ("path", "text", "rect", "circle", "ellipse",
                                  "line", "polygon", "polyline")
               for el in probe.iter()):
        return None
    try:
        img = render_svg_to_png(stripped, width=int(vbw * overlay_scale))
    except Exception as e:
        logger.warning(f"SVG overlay render failed, skipping decorations: {e}")
        return None
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# --------------------------------------------------------------------------- #
# public API
# --------------------------------------------------------------------------- #
class SvgToPptxService:
    """Translate a deck's SVG pages into one native, editable PPTX."""

    @staticmethod
    def _read_viewbox(root) -> Tuple[float, float]:
        vb = root.get("viewBox")
        if vb:
            parts = re.split(r"[ ,]+", vb.strip())
            if len(parts) == 4:
                try:
                    return float(parts[2]), float(parts[3])
                except ValueError:
                    pass
        # fall back to width/height or the contract default
        try:
            return float(root.get("width", DEFAULT_VBW)), float(root.get("height", DEFAULT_VBH))
        except (ValueError, TypeError):
            return DEFAULT_VBW, DEFAULT_VBH

    @staticmethod
    def _build_slide_native(prs, svg_text, stats):
        """Translate every primitive to a native pptx shape + one decorative overlay."""
        root = etree.fromstring(svg_text.encode("utf-8"))
        vbw, vbh = SvgToPptxService._read_viewbox(root)
        scale = _Scale(vbw, vbh)
        shapes = prs.slides.add_slide(prs.slide_layouts[6]).shapes

        for el in root.iter():
            name = _localname(el)
            if name not in SHAPELIKE or _in_faint_group(el):
                continue
            try:
                if name == "rect":
                    _add_rect(shapes, el, root, scale); stats["shapes"] += 1
                elif name in ("circle", "ellipse"):
                    _add_oval(shapes, el, scale); stats["shapes"] += 1
                elif name == "line":
                    _add_line(shapes, el, scale); stats["lines"] += 1
                elif name == "text":
                    _add_text(shapes, el, scale); stats["text"] += 1
            except Exception as e:
                logger.warning(f"Skipped <{name}> during SVG→PPTX translation: {e}")

        overlay = _render_overlay(svg_text, vbw, 2)
        if overlay is not None:
            shapes.add_picture(overlay, 0, 0, prs.slide_width, prs.slide_height)
            stats["overlays"] += 1

    @staticmethod
    def _build_slide_raster(prs, svg_text, stats, bg_width):
        """Pixel-faithful background (resvg) + native, editable text boxes on top."""
        root = etree.fromstring(svg_text.encode("utf-8"))
        vbw, vbh = SvgToPptxService._read_viewbox(root)
        scale = _Scale(vbw, vbh)
        shapes = prs.slides.add_slide(prs.slide_layouts[6]).shapes

        bg = _render_background(svg_text, bg_width)
        if bg is not None:
            shapes.add_picture(bg, 0, 0, prs.slide_width, prs.slide_height)
            stats["backgrounds"] += 1

        # Re-add the editable text as native boxes, positioned to match the SVG.
        for el in _editable_texts(root):
            try:
                _add_text(shapes, el, scale); stats["text"] += 1
            except Exception as e:
                logger.warning(f"Skipped <text> during SVG→PPTX translation: {e}")

    @staticmethod
    def _collect_text_chars(svgs: List[str]) -> str:
        """Every character that appears in a <text> across the deck (for font subsetting)."""
        chars = set()
        for svg_text in svgs:
            try:
                root = etree.fromstring(svg_text.encode("utf-8"))
            except Exception:
                continue
            for el in root.iter():
                if _localname(el) == "text":
                    for s in el.itertext():
                        chars.update(s)
        return "".join(sorted(chars))

    @staticmethod
    def build_from_strings(svgs: List[str], output_file: Optional[str] = None,
                           mode: str = "raster", bg_width: int = 2560,
                           embed_font: bool = True
                           ) -> Tuple[Optional[bytes], dict]:
        """Build a PPTX from a list of SVG document strings.

        Args:
            svgs: SVG documents, one per slide, in order.
            output_file: where to write the .pptx; if None, returns bytes.
            mode: 'raster' (default) = pixel-faithful resvg background + editable text
                  boxes — matches the web preview exactly except for text font; 'native'
                  = translate every primitive to an editable pptx shape (more editable,
                  but PowerPoint can't reproduce SVG gradients/opacity/paths 1:1).
            bg_width: background raster width in px (raster mode only).

        Returns:
            (pptx_bytes_or_None, stats).
        """
        if not svgs:
            raise ValueError("No SVG content provided")
        mode = (mode or "raster").lower()
        if mode not in ("raster", "native"):
            raise ValueError(f"Unknown mode: {mode!r} (expected 'raster' or 'native')")

        prs = Presentation()
        prs.slide_width = Emu(SLIDE_W_EMU)
        prs.slide_height = Emu(SLIDE_H_EMU)

        stats = {"slides": 0, "shapes": 0, "text": 0, "lines": 0,
                 "overlays": 0, "backgrounds": 0, "mode": mode}

        for svg_text in svgs:
            if mode == "raster":
                SvgToPptxService._build_slide_raster(prs, svg_text, stats, bg_width)
            else:
                SvgToPptxService._build_slide_native(prs, svg_text, stats)
            stats["slides"] += 1

        buf = BytesIO()
        prs.save(buf)
        data = buf.getvalue()

        # Embed the (subsetted) font so the deck's text looks identical to the SVG
        # preview on any machine — no reliance on the viewer having Noto installed.
        stats["font_embedded"] = False
        if embed_font:
            chars = SvgToPptxService._collect_text_chars(svgs)
            if chars:
                from services.pptx_font_embed import embed_font_into_pptx_bytes
                from services.svg_render_service import _FONTS_DIR
                font_path = os.path.join(_FONTS_DIR, "NotoSansSC-Regular.ttf")
                new_data = embed_font_into_pptx_bytes(data, chars, DEFAULT_FONT_FAMILY, font_path)
                stats["font_embedded"] = new_data is not data and len(new_data) != len(data)
                data = new_data

        if output_file:
            with open(output_file, "wb") as f:
                f.write(data)
            return None, stats
        return data, stats

    @staticmethod
    def build_from_paths(svg_paths: List[str], output_file: Optional[str] = None,
                         mode: str = "raster", bg_width: int = 2560,
                         embed_font: bool = True
                         ) -> Tuple[Optional[bytes], dict]:
        """Build a PPTX from a list of .svg file paths (in order)."""
        svgs = []
        for p in svg_paths:
            if not os.path.exists(p):
                logger.warning(f"SVG file not found, skipping: {p}")
                continue
            with open(p, "r", encoding="utf-8") as f:
                svgs.append(f.read())
        if not svgs:
            raise ValueError("No readable SVG files found")
        return SvgToPptxService.build_from_strings(svgs, output_file, mode, bg_width, embed_font)
